from pptx import Presentation

fichier_pptx = "C:/Users/hp/Documents/Stage/sonatel_iot_journey/Bilan_Pilote_bouton_de_satisfaction_LoRa_v16092021_DG.pptx"
presentation = Presentation(fichier_pptx)

with open("extraction8.txt", "w", encoding="utf-8") as f:
    for i, slide in enumerate(presentation.slides, start=1):
        f.write(f"\n--- Slide {i} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texte = shape.text.strip()
                if texte:
                    f.write(texte + "\n")

print("Texte extrait avec succès dans extraction.txt")
