
from pptx import Presentation
import os

def extraire_images_de_pptx(pptx_path, output_base_dir):
    filename = os.path.splitext(os.path.basename(pptx_path))[0]
    output_dir = os.path.join(output_base_dir, filename)
    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation(pptx_path)

    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # 13 = image
                try:
                    image = shape.image
                    extension = image.ext
                    image_bytes = image.blob
                    image_filename = f"slide_{i+1}_img_{j+1}.{extension}"
                    image_path = os.path.join(output_dir, image_filename)

                    with open(image_path, "wb") as f:
                        f.write(image_bytes)

                    print(f"[✔] {image_filename} extrait de {filename}")
                except Exception as e:
                    print(f"[⚠] Erreur dans {filename}, slide {i+1} : {e}")

def traiter_tous_les_pptx(dossier_pptx, dossier_images):
    for fichier in os.listdir(dossier_pptx):
        if fichier.endswith(".pptx"):
            chemin_fichier = os.path.join(dossier_pptx, fichier)
            extraire_images_de_pptx(chemin_fichier, dossier_images)

# === Utilisation ===
dossier_pptx = "C:/Users/hp/Documents/Stage/sonatel_iot_journey"
dossier_images = os.path.join(dossier_pptx, "images_extraites")

traiter_tous_les_pptx(dossier_pptx, dossier_images)
